"""LLM-driven PowerPoint (.pptx) builder.

The LLM authors each slide as raw PresentationML/DrawingML XML — a
<p:txBody> for a placeholder's text, one or more <p:sp> shapes, or the full
shape list of a slide — and the `create_pptx` tool hands them to
PptxBuilder, which injects them into an unpacked .pptx package.

Injection uses defusedxml.minidom so namespace prefixes are preserved exactly
as written (ElementTree would rewrite them and corrupt the file on save).
This mirrors the reference slide-injector utility's `replace_txbody` /
`replace_whole_slide_body` approach.
"""

import os
import re
import shutil
import tempfile
import zipfile
from xml.sax.saxutils import escape as _escape

from defusedxml import minidom

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# Wrapping the fragment in a root that declares every namespace the LLM might
# reference lets a bare fragment (no xmlns of its own) parse standalone,
# exactly as if it were still nested inside <p:sld>.
_NS_WRAPPER = '<root xmlns:a="{a}" xmlns:p="{p}" xmlns:r="{r}">{fragment}</root>'

_XML_DECL = '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'


def parse_fragment(xml_text: str):
    """Parse a namespace-qualified XML fragment (e.g. from an LLM) that has
    no xmlns declarations of its own and return its single root element."""
    wrapped = _NS_WRAPPER.format(a=A_NS, p=P_NS, r=R_NS, fragment=xml_text.strip())
    frag_doc = minidom.parseString(wrapped)
    root = frag_doc.documentElement
    children = [n for n in root.childNodes if n.nodeType == n.ELEMENT_NODE]
    if len(children) != 1:
        raise ValueError(f"Expected exactly one root element in fragment, got {len(children)}")
    return children[0]


def parse_fragments(xml_text: str):
    """Parse a fragment that may contain several top-level elements (e.g. a
    list of <p:sp> shapes) and return them in document order."""
    wrapped = _NS_WRAPPER.format(a=A_NS, p=P_NS, r=R_NS, fragment=xml_text.strip())
    frag_doc = minidom.parseString(wrapped)
    root = frag_doc.documentElement
    return [n for n in root.childNodes if n.nodeType == n.ELEMENT_NODE]


# ── plain-text -> DrawingML helpers (for title/body without raw XML) ──────


def _para(text: str, size: int = 1400, bold: bool = False, bullet: bool = False) -> str:
    if bullet:
        ppr = '<a:pPr marL="228600" indent="-228600"><a:buChar char="•"/></a:pPr>'
    else:
        ppr = '<a:pPr marL="0" indent="0"><a:buNone/></a:pPr>'
    b = ' b="1"' if bold else ""
    return (
        f'<a:p>{ppr}<a:r><a:rPr lang="en-US" dirty="0" sz="{size}"{b}/>'
        f"<a:t>{_escape(text)}</a:t></a:r></a:p>"
    )


def _txbody(paragraphs) -> str:
    return f'<p:txBody><a:bodyPr wrap="square"/><a:lstStyle/>{"".join(paragraphs)}</p:txBody>'


def title_txbody(text: str) -> str:
    return _txbody([_para(str(text), size=3200, bold=True)])


def body_txbody(lines) -> str:
    if isinstance(lines, str):
        lines = lines.split("\n")
    paras = []
    for ln in lines:
        ln = str(ln).strip()
        if not ln:
            continue
        bullet = ln[0] in "-*"
        if bullet:
            ln = ln[1:].strip()
        paras.append(_para(ln, size=1400, bullet=bullet))
    if not paras:
        paras.append(_para(""))
    return _txbody(paras)


class PptxBuilder:
    """Build an editable .pptx by injecting LLM-authored XML fragments into a
    python-pptx-generated template package (or a user-provided template)."""

    def __init__(self, num_slides: int = 1, template: str | None = None):
        self._dir = tempfile.mkdtemp(prefix="pptx_builder_")
        self._slides_dir = os.path.join(self._dir, "ppt", "slides")
        if template is not None:
            self._unpack(template, self._dir)
        else:
            seed = os.path.join(self._dir, "_seed.pptx")
            from pptx import Presentation
            prs = Presentation()
            layout = prs.slide_layouts[1]  # Title and Content
            for _ in range(max(1, int(num_slides))):
                prs.slides.add_slide(layout)
            prs.save(seed)
            self._unpack(seed, self._dir)
            os.remove(seed)
        self._slide_files = self._list_slide_files()

    # ── package plumbing ────────────────────────────────────────────────

    @staticmethod
    def _unpack(src: str, dst: str) -> None:
        with zipfile.ZipFile(src) as zf:
            zf.extractall(dst)

    def _list_slide_files(self):
        names = [f for f in os.listdir(self._slides_dir) if re.fullmatch(r"slide\d+\.xml", f)]
        names.sort(key=lambda f: int(re.sub(r"\D", "", f)))
        return names

    def save(self, path: str) -> str:
        """Repack the working directory as a .pptx and write it to `path`."""
        path = os.path.abspath(path)
        os.makedirs(os.path.dirname(path), exist_ok=True)
        entries = []
        for root, _dirs, files in os.walk(self._dir):
            for fn in files:
                entries.append(os.path.relpath(os.path.join(root, fn), self._dir))
        # [Content_Types].xml must be the first entry in the zip.
        entries.sort(key=lambda n: (n != "[Content_Types].xml", n))
        with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
            for name in entries:
                zf.write(os.path.join(self._dir, name), arcname=name)
        return path

    def cleanup(self) -> None:
        shutil.rmtree(self._dir, ignore_errors=True)

    # ── slide XML plumbing ───────────────────────────────────────────────

    def _slide_path(self, slide_idx: int) -> str:
        idx = int(slide_idx)
        if not 1 <= idx <= len(self._slide_files):
            raise ValueError(
                f"slide index {slide_idx} out of range (deck has {len(self._slide_files)} slides)"
            )
        return os.path.join(self._slides_dir, self._slide_files[idx - 1])

    @staticmethod
    def _save(doc, path: str) -> None:
        xml_out = doc.toxml()
        body = xml_out[xml_out.find("?>") + 2:]  # drop the decl minidom emits
        with open(path, "w", encoding="utf-8") as f:
            f.write(_XML_DECL + "\n" + body)

    @staticmethod
    def _find_placeholder_sp(doc, placeholder_type: str | None = None, placeholder_idx: str | None = None):
        """Locate a placeholder <p:sp> by ph type or idx. If neither is given,
        the first non-title placeholder (the body/content box) is used."""
        for sp in doc.getElementsByTagNameNS(P_NS, "sp"):
            ph_list = sp.getElementsByTagNameNS(P_NS, "ph")
            if not ph_list:
                continue
            ph = ph_list[0]
            if placeholder_idx is not None:
                if ph.getAttribute("idx") == str(placeholder_idx):
                    return sp
            elif placeholder_type is not None:
                if ph.getAttribute("type") == str(placeholder_type):
                    return sp
            elif ph.getAttribute("type") != "title":
                return sp
        return None

    # ── mutations ────────────────────────────────────────────────────────

    def set_placeholder_xml(self, slide_idx: int, xml_text: str,
                            placeholder_type: str | None = None,
                            placeholder_idx: str | None = None) -> None:
        """Replace a placeholder's <p:txBody> with an LLM-authored
        <p:txBody> fragment. Target by ph type ('title', ...) or idx (the
        body/content box is typically idx="1"); if neither is given, the
        first non-title placeholder is used."""
        path = self._slide_path(slide_idx)
        doc = minidom.parse(path)
        sp = self._find_placeholder_sp(doc, placeholder_type, placeholder_idx)
        if sp is None:
            raise ValueError(
                f"No placeholder found (type={placeholder_type!r}, idx={placeholder_idx!r}) in slide {slide_idx}"
            )
        new_txbody = parse_fragment(xml_text)
        if new_txbody.localName != "txBody":
            raise ValueError(
                f"txbody mode expects a <p:txBody> fragment, got <{new_txbody.nodeName}>"
            )
        old_txbody = sp.getElementsByTagNameNS(P_NS, "txBody")[0]
        sp.replaceChild(doc.importNode(new_txbody, deep=True), old_txbody)
        self._save(doc, path)

    def set_slide_shapes_xml(self, slide_idx: int, xml_text: str) -> None:
        """Replace everything inside <p:spTree> with LLM-authored shapes —
        use when the LLM is authoring the full slide layout, not just one
        placeholder's text. The mandatory <p:nvGrpSpPr>/<p:grpSpPr>
        housekeeping nodes are kept; all placeholders/shapes are removed."""
        path = self._slide_path(slide_idx)
        doc = minidom.parse(path)
        sp_tree = doc.getElementsByTagNameNS(P_NS, "spTree")[0]
        keep = {"nvGrpSpPr", "grpSpPr"}
        for child in list(sp_tree.childNodes):
            if child.nodeType == child.ELEMENT_NODE and child.localName not in keep:
                sp_tree.removeChild(child)
        shapes = parse_fragments(xml_text)
        if not shapes:
            raise ValueError("spTree fragment contained no shape elements")
        for el in shapes:
            sp_tree.appendChild(doc.importNode(el, deep=True))
        self._save(doc, path)

    def append_shapes_xml(self, slide_idx: int, xml_text: str) -> None:
        """Append LLM-authored shapes to <p:spTree> after the existing
        placeholders — use to add charts/boxes while keeping title/body."""
        path = self._slide_path(slide_idx)
        doc = minidom.parse(path)
        sp_tree = doc.getElementsByTagNameNS(P_NS, "spTree")[0]
        shapes = parse_fragments(xml_text)
        if not shapes:
            raise ValueError("append fragment contained no shape elements")
        for el in shapes:
            sp_tree.appendChild(doc.importNode(el, deep=True))
        self._save(doc, path)

    def set_title(self, slide_idx: int, text: str) -> None:
        self.set_placeholder_xml(slide_idx, title_txbody(text), placeholder_type="title")

    def set_body_text(self, slide_idx: int, lines) -> None:
        self.set_placeholder_xml(slide_idx, body_txbody(lines))
