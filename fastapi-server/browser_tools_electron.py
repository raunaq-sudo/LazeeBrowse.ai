import json
import os
import re
import datetime
from typing import List
from langchain.tools import tool


def extract_text_from_file(full_path: str) -> str:
    """Extract readable text from a file based on extension.

    Supports PDF, DOCX, XLSX, PPTX and falls back to plain UTF-8 text for
    everything else. On parse errors returns a short error note instead of
    raising, so the agent can react gracefully.
    """
    ext = os.path.splitext(full_path)[1].lower()
    try:
        if ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(full_path)
            return "\n".join((page.extract_text() or "") for page in reader.pages)
        if ext == ".docx":
            from docx import Document
            doc = Document(full_path)
            parts = [p.text for p in doc.paragraphs if p.text]
            for table in doc.tables:
                for row in table.rows:
                    parts.append(" | ".join(c.text.strip() for c in row.cells))
            return "\n".join(parts)
        if ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(full_path, data_only=True, read_only=True)
            parts = []
            for ws in wb.worksheets:
                parts.append(f"[Sheet: {ws.title}]")
                for row in ws.iter_rows(values_only=True):
                    vals = ["" if v is None else str(v) for v in row]
                    if any(vals):
                        parts.append(" | ".join(vals))
            return "\n".join(parts)
        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(full_path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"[Slide {i}]")
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        for para in shape.text_frame.paragraphs:
                            txt = "".join(run.text for run in para.runs)
                            if txt.strip():
                                parts.append(txt)
                    if getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            parts.append(" | ".join(c.text.strip() for c in row.cells))
            return "\n".join(parts)
        with open(full_path, "r", encoding="utf-8", errors="replace") as fh:
            return fh.read()
    except Exception as e:
        return f"[Error reading {os.path.basename(full_path)}: {e}]"


_TABLE_SEP = re.compile(r"^[\|\s:\-]+$")


def _html_escape(text: str) -> str:
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _write_docx(path: str, content: str) -> None:
    from docx import Document
    doc = Document()
    lines = content.split("\n")
    i = 0
    while i < len(lines):
        line = lines[i].rstrip()
        if line.startswith("|") and i + 1 < len(lines) and _TABLE_SEP.match(lines[i + 1].strip()):
            header = [c.strip() for c in line.strip("|").split("|")]
            table = doc.add_table(rows=1, cols=max(1, len(header)))
            table.style = "Table Grid"
            for ci, h in enumerate(header):
                table.rows[0].cells[ci].text = h
            i += 2
            while i < len(lines) and lines[i].strip().startswith("|"):
                row = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                cells = table.add_row().cells
                for ci in range(min(len(row), len(cells))):
                    cells[ci].text = row[ci]
                i += 1
            continue
        m = re.match(r"^(#{1,6})\s+(.*)", line)
        if m:
            doc.add_heading(m.group(2), level=min(len(m.group(1)), 4))
        elif line.strip():
            doc.add_paragraph(line)
        i += 1
    doc.save(path)


def _write_xlsx(path: str, content: str) -> None:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    for line in content.split("\n"):
        line = line.rstrip("\r")
        if not line.strip():
            continue
        if "\t" in line:
            row = [c.strip() for c in line.split("\t")]
        else:
            row = [c.strip() for c in line.split(",")]
        ws.append(row)
    wb.save(path)


def _write_pdf(path: str, content: str, landscape: bool = False, slides: bool = False) -> None:
    from reportlab.lib.pagesizes import letter, landscape as _landscape
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    styles = getSampleStyleSheet()
    pagesize = _landscape(letter) if landscape else letter
    margin = 0.55 * inch if landscape else 0.9 * inch
    doc = SimpleDocTemplate(path, pagesize=pagesize, margins=(margin, margin, margin, margin))
    story = []
    blocks = re.split(r"\n\s*-{3,}\s*\n", content) if slides else [content]
    for bi, block in enumerate(blocks):
        if slides and bi > 0:
            story.append(PageBreak())
        for line in block.split("\n"):
            line = line.rstrip()
            if not line.strip():
                story.append(Spacer(1, 6))
                continue
            m = re.match(r"^(#{1,3})\s+(.*)", line)
            if m:
                lvl = len(m.group(1))
                style = styles["Heading1"] if lvl == 1 else (styles["Heading2"] if lvl == 2 else styles["Heading3"])
                story.append(Paragraph(_html_escape(m.group(2)), style))
            else:
                story.append(Paragraph(_html_escape(line), styles["BodyText"]))
    doc.build(story)


def _looks_like_html(content: str) -> bool:
    s = content.lstrip().lower()
    return s.startswith((
        "<!doctype", "<html", "<head", "<body", "<style", "<section",
        "<div", "<h1", "<h2", "<h3", "<ul", "<ol", "<table", "<p",
    ))


async def _write_pdf_html(path: str, content: str, width: str | None = None, height: str | None = None) -> None:
    """Render HTML+CSS to PDF using headless Chromium (full CSS support)."""
    from playwright.async_api import async_playwright
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        try:
            page = await browser.new_page()
            await page.set_content(content, wait_until="load")
            pdf_kwargs = {"print_background": True}
            if width and height:
                pdf_kwargs.update(width=width, height=height,
                                  margin={"top": "0", "bottom": "0", "left": "0", "right": "0"})
            else:
                pdf_kwargs["format"] = "A4"
            await page.pdf(path=path, **pdf_kwargs)
        finally:
            await browser.close()


# ── PPTX (native, editable) ─────────────────────
# Slides are built directly with python-pptx using real text frames/runs so
# the output stays editable. HTML+CSS goes through a lightweight CSS
# resolver; markdown is split into one slide per heading (configurable level).

_CSS_STYLE_PROPS = ("color", "font-size", "font-weight", "font-family", "text-align", "background-color")


def _style_dict(style_str: str) -> dict:
    d = {}
    if not style_str:
        return d
    for part in style_str.split(";"):
        if ":" in part:
            k, v = part.split(":", 1)
            d[k.strip().lower()] = v.strip()
    return d


def _css_size_to_pt(value: str) -> int | None:
    m = re.match(r"^([\d.]+)\s*(px|pt)?$", value.strip().lower())
    if not m:
        return None
    num = float(m.group(1))
    unit = m.group(2)
    if unit == "px":
        num = num * 0.75
    return max(8, int(num))


def _css_color_to_hex(color: str) -> str | None:
    """Normalize a CSS color value to a 6-digit hex string, or None if unsupported."""
    color = (color or "").strip()
    if color.startswith("#"):
        hex_color = color.lstrip("#")
        if len(hex_color) == 3:
            hex_color = "".join(c * 2 for c in hex_color)
        return hex_color if re.fullmatch(r"[0-9a-fA-F]{6}", hex_color) else None
    named = {"red": "FF0000", "black": "000000", "white": "FFFFFF", "blue": "0000FF",
             "green": "008000", "gray": "808080", "grey": "808080", "yellow": "FFFF00"}
    return named.get(color.lower())


def _parse_css_rules(css_text: str) -> List:
    css_text = re.sub(r"/\*.*?\*/", "", css_text, flags=re.S)
    rules = []
    for sel_group, body in re.findall(r"([^{}]+)\{([^{}]*)\}", css_text):
        declarations = {}
        for part in body.split(";"):
            if ":" in part:
                k, v = part.split(":", 1)
                if k.strip().lower() in _CSS_STYLE_PROPS:
                    declarations[k.strip().lower()] = v.strip()
        if declarations:
            for sel in sel_group.split(","):
                sel = sel.strip()
                if sel:
                    rules.append((sel, declarations))
    return rules


def _css_specificity(sel: str) -> tuple:
    ids = len(re.findall(r"#[a-zA-Z0-9_-]+", sel))
    classes = len(re.findall(r"\.[a-zA-Z0-9_-]+", sel))
    tags = len(re.findall(r"(?<![.#])\b[a-zA-Z0-9]+", sel))
    return (ids, classes, tags)


def _simple_css_match(sel: str, el) -> bool:
    sel = sel.strip()
    if sel.startswith("."):
        return sel[1:] in (el.get("class") or "").split()
    if sel.startswith("#"):
        return el.get("id") == sel[1:]
    m = re.match(r"^([a-zA-Z][a-zA-Z0-9]*)((?:\.[a-zA-Z0-9_-]+)*)$", sel)
    if not m:
        return False
    tag, classes = m.group(1), m.group(2)
    if el.tag != tag:
        return False
    el_classes = (el.get("class") or "").split()
    for cls in re.findall(r"\.([a-zA-Z0-9_-]+)", classes):
        if cls not in el_classes:
            return False
    return True


def _css_selector_matches(sel: str, el, root) -> bool:
    parts = [p for p in sel.split() if p]
    if not parts:
        return False
    if not _simple_css_match(parts[-1], el):
        return False
    for ancestor in parts[-2::-1]:
        el = el.getparent()
        while el is not None and el is not root:
            if _simple_css_match(ancestor, el):
                break
            el = el.getparent()
        if el is None or el is root:
            return False
    return True


def _resolve_element_style(el, rules: List, root) -> dict:
    """Resolve the effective style for a single element (specificity + source
    order for ties, then inline styles last)."""
    eff = {}
    for i, (sel, decls) in enumerate(rules):
        if not _css_selector_matches(sel, el, root):
            continue
        spec = _css_specificity(sel)
        for k, v in decls.items():
            cur = eff.get(k)
            if cur is None or (spec, i) >= (cur[0], cur[2]):
                eff[k] = (spec, v, i)
    style = {k: v for k, (_, v, _) in eff.items()}
    for k, v in _style_dict(el.get("style") or "").items():
        if k in _CSS_STYLE_PROPS:
            style[k] = v
    return style


def _write_pptx_html(path: str, content: str) -> None:
    """Convert HTML+CSS into editable slides. Elements with class 'slide'
    (or <section>) become slides; the first <h1> is the title, other blocks
    the body. Stylesheet rules (tag/.class/#id/descendant selectors) are
    resolved and applied to color/font-size/font-weight/font-family/text-align."""
    from lxml import html as lxml_html
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    root = lxml_html.fromstring(content)
    style_text = "".join((s.text or "") for s in root.xpath("//style"))
    rules = _parse_css_rules(style_text)

    slides_elems = root.xpath(".//*[contains(concat(' ', normalize-space(@class), ' '), ' slide ')]")
    if not slides_elems:
        slides_elems = root.xpath(".//section")
    if not slides_elems:
        slides_elems = [root]

    def _align(style):
        align = style.get("text-align")
        return {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "left": PP_ALIGN.LEFT}.get(align)

    def _apply_run(run, style, kind):
        color = style.get("color")
        size = style.get("font-size")
        weight = style.get("font-weight")
        family = style.get("font-family")
        if weight in ("bold", "bolder") or (weight and weight.isdigit() and int(weight) >= 600):
            run.font.bold = True
        elif weight in ("normal", "400"):
            run.font.bold = False
        elif kind == "heading":
            run.font.bold = True
        if color:
            hex_color = _css_color_to_hex(color)
            if hex_color:
                try:
                    run.font.color.rgb = RGBColor.from_string(hex_color)
                except Exception:
                    pass
        if size:
            pt = _css_size_to_pt(size)
            if pt:
                run.font.size = Pt(pt)
        if family:
            run.font.name = family.split(",")[0].strip().strip("'\"")

    def _apply_background(slide, style):
        bg = style.get("background-color")
        hex_color = _css_color_to_hex(bg) if bg else None
        if hex_color:
            try:
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor.from_string(hex_color)
            except Exception:
                pass

    prs = Presentation()
    layout = prs.slide_layouts[1]
    _BLOCK_TAGS = {"h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "div", "section", "ul", "ol", "table", "tr", "td", "th"}
    for se in slides_elems:
        title = ""
        title_style = {}
        body_items = []
        for el in se.iter():
            if el is not se and (el.tag in ("section", "div")) and "slide" in (el.get("class") or "").split():
                continue
            if el.tag in ("div", "section") and any(c.tag in _BLOCK_TAGS for c in el):
                continue
            if el.tag in ("h1", "h2", "h3", "h4", "p", "li", "div", "td", "th"):
                text = (el.text_content() or "").strip()
                if not text:
                    continue
                style = _resolve_element_style(el, rules, root)
                if el.tag == "h1" and not title:
                    title = text
                    title_style = style
                elif el.tag == "li":
                    body_items.append(("bullet", text, style))
                elif el.tag in ("h2", "h3", "h4"):
                    body_items.append(("heading", text, style))
                else:
                    body_items.append(("para", text, style))
        slide = prs.slides.add_slide(layout)
        bg_style = _resolve_element_style(se, rules, root)
        if not bg_style.get("background-color"):
            bg_style = _resolve_element_style(root, rules, root)
        _apply_background(slide, bg_style)
        if title:
            title_tf = slide.shapes.title.text_frame
            title_tf.clear()
            title_tf.text = title
            for para in title_tf.paragraphs:
                for run in para.runs:
                    _apply_run(run, title_style, "heading")
        body_shape = None
        for ph in slide.placeholders:
            if ph.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                body_shape = ph
                break
        if body_items and body_shape is not None:
            tf = body_shape.text_frame
            tf.word_wrap = True
            tf.clear()
            first = True
            for kind, text, style in body_items:
                para = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                para.text = ("• " + text) if kind == "bullet" else text
                if _align(style) is not None:
                    para.alignment = _align(style)
                for run in para.runs:
                    _apply_run(run, style, kind)
    prs.save(path)


def _split_markdown_into_slides(content: str, heading_level: int = 1):
    """Split markdown into slides at every heading of `heading_level` or
    shallower (default: every H1). Returns [(title, body_lines)]."""
    heading_re = re.compile(r"^(#{1,6})\s+(.*)")
    slides = []
    current_title, current_body = None, []

    def flush():
        if current_title is not None or any(b.strip() for b in current_body):
            slides.append((current_title or "", current_body))

    for line in content.split("\n"):
        m = heading_re.match(line.rstrip())
        if m and len(m.group(1)) <= heading_level:
            flush()
            current_title, current_body = m.group(2).strip(), []
        elif line.strip() == "---":
            continue  # horizontal rules are separators, not slide content
        else:
            current_body.append(line)
    flush()
    if not slides:
        slides = [("", content.split("\n"))]
    return slides


_SLIDE_LEVEL_RE = re.compile(r"^\s*#\s*slide[-_]?level\s*[:=]\s*([1-6])\s*$", re.I)
_HTML_SLIDE_LEVEL_RE = re.compile(r"^\s*<!--\s*slide[-_]?level\s*[:=]\s*([1-6])\s*-->\s*$", re.I)


def _extract_slide_level(content: str):
    """Read an optional 'slide-level: N' directive from the top of the content.
    Returns (heading_level, content_without_directive)."""
    lines = content.split("\n")
    for i, line in enumerate(lines):
        m = _SLIDE_LEVEL_RE.match(line) or _HTML_SLIDE_LEVEL_RE.match(line)
        if m:
            del lines[i]
            return int(m.group(1)), "\n".join(lines)
    return 1, content


def _write_pptx_native(path: str, content: str) -> None:
    """Build a real, editable .pptx. Accepts HTML+CSS (via the CSS-aware
    _write_pptx_html path) or markdown — markdown is split into one slide per
    heading so decks don't need explicit slide markers."""
    if _looks_like_html(content):
        _write_pptx_html(path, content)
        return

    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER

    heading_level, content = _extract_slide_level(content)
    prs = Presentation()
    layout = prs.slide_layouts[1]
    for title, body_lines in _split_markdown_into_slides(content, heading_level=heading_level):
        slide = prs.slides.add_slide(layout)
        if title:
            slide.shapes.title.text_frame.text = title

        body_shape = next(
            (ph for ph in slide.placeholders
             if ph.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)),
            None,
        )
        if body_shape is None:
            continue
        tf = body_shape.text_frame
        tf.word_wrap = True
        tf.clear()
        first = True
        for line in body_lines:
            line = line.strip()
            if not line:
                continue
            sub = re.match(r"^(#{2,6})\s+(.*)", line)  # sub-heading inside a slide
            bullet = re.match(r"^[-*]\s+(.*)", line)
            text = sub.group(2) if sub else (bullet.group(1) if bullet else line)
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            para.text = ("• " + text) if bullet else text
            if sub:
                for run in para.runs:
                    run.font.bold = True
    prs.save(path)


async def write_text_to_file(full_path: str, content: str, append: bool = False) -> str:
    """Write content to a file, converting to the target format by extension."""
    ext = os.path.splitext(full_path)[1].lower()
    if ext in (".docx", ".xlsx", ".pptx", ".pdf"):
        if append:
            return "ERROR: append is not supported for this file type"
        try:
            if ext == ".docx":
                _write_docx(full_path, content)
            elif ext == ".xlsx":
                _write_xlsx(full_path, content)
            elif ext == ".pptx":
                _write_pptx_native(full_path, content)
                return f"OK: {os.path.basename(full_path)}"
            elif ext == ".pdf":
                if _looks_like_html(content):
                    await _write_pdf_html(full_path, content)
                else:
                    _write_pdf(full_path, content)
            return f"OK: {os.path.basename(full_path)}"
        except Exception as e:
            return f"ERROR: {e}"
    mode = "a" if append else "w"
    with open(full_path, mode, encoding="utf-8") as f:
        f.write(content)
    return f"OK: {os.path.basename(full_path)}"


def build_tools(browser_command, log_chat, base_path=None):
    """Create LangChain tools that control Electron's BrowserView via IPC."""

    def get_user_files_dir():
        base = os.path.join(base_path, "files") if base_path else "files"
        os.makedirs(base, exist_ok=True)
        return base

    def resolve_user_path(relative_path: str):
        base = get_user_files_dir()
        clean = os.path.normpath(relative_path).lstrip(os.sep)
        full_path = os.path.join(base, clean)
        if not os.path.abspath(full_path).startswith(os.path.abspath(base)):
            raise Exception("Invalid file path (security violation)")
        return full_path

    def refresh_file_tree():
        pass

    # ── BROWSER: NAVIGATION ──────────────────────

    @tool
    async def open_url(url: str) -> str:
        """Navigate the browser to a URL."""
        await log_chat(f"Navigating to {url}")
        result = await browser_command("navigate", {"url": url})
        if isinstance(result, dict) and result.get("error"):
            return f"Error: {result['error']}"
        return f"Opened {url}"

    @tool
    async def get_url() -> str:
        """Get the current page URL."""
        return await browser_command("get_url", {})

    @tool
    async def get_title() -> str:
        """Get the current page title."""
        return await browser_command("get_title", {})

    @tool
    async def go_back() -> str:
        """Navigate back."""
        await browser_command("go_back", {})
        return "OK"

    @tool
    async def go_forward() -> str:
        """Navigate forward."""
        await browser_command("go_forward", {})
        return "OK"

    # ── BROWSER: INTERACTION ─────────────────────

    @tool
    async def click(selector: str) -> str:
        """Click an element by CSS selector."""
        await log_chat(f"Clicking {selector}")
        result = await browser_command("click", {"selector": selector})
        if isinstance(result, dict) and result.get("error"):
            return f"Error: {result['error']}"
        return f"Clicked {selector}"

    @tool
    async def type_text(selector: str, text: str) -> str:
        """Type text into an input field."""
        await log_chat(f"Typing into {selector}")
        result = await browser_command("type", {"selector": selector, "text": text})
        if isinstance(result, dict) and result.get("error"):
            return f"Error: {result['error']}"
        return f"Typed into {selector}"

    @tool
    async def scroll(amount: int = 500) -> str:
        """Scroll down (positive) or up (negative) by N pixels."""
        await log_chat(f"Scrolling {amount}px")
        result = await browser_command("scroll", {"amount": amount})
        if isinstance(result, dict) and result.get("error"):
            return f"Error: {result['error']}"
        if isinstance(result, dict):
            return f"Scrolled {amount}px (delta: {result.get('delta', '?')}, after: {result.get('after', '?')})"
        return f"Scrolled {amount}px"

    @tool
    async def submit_form() -> str:
        """Submit the current form."""
        await browser_command("submit_form", {})
        return "Form submitted"

    @tool
    async def press_key(key: str) -> str:
        """Press a keyboard key (Enter, Tab, Escape, ArrowDown, etc.)."""
        await browser_command("press_key", {"key": key})
        return f"Pressed {key}"

    @tool
    async def page_down() -> str:
        """Scroll down one page (PageDown key)."""
        await browser_command("press_key", {"key": "PageDown"})
        return "Page down"

    @tool
    async def page_up() -> str:
        """Scroll up one page (PageUp key)."""
        await browser_command("press_key", {"key": "PageUp"})
        return "Page up"

    @tool
    async def select_option(selector: str, value: str = "", label: str = "") -> str:
        """Select an option in a <select> dropdown by value or label text."""
        await log_chat(f"Selecting option in {selector}")
        result = await browser_command("select_option", {"selector": selector, "value": value, "label": label})
        if isinstance(result, dict) and result.get("error"):
            return f"Error: {result['error']}"
        return f"Selected option in {selector}"

    @tool
    async def get_dropdown_options(selector: str) -> list:
        """Get all options from a <select> dropdown. Returns [{value, label, selected}]."""
        await log_chat(f"Getting dropdown options for {selector}")
        return await browser_command("get_dropdown_options", {"selector": selector})

    # ── BROWSER: EXTRACTION ──────────────────────

    @tool
    async def get_page_text() -> str:
        """Get all visible text from the page."""
        await log_chat("Getting page text")
        return await browser_command("get_text", {})

    @tool
    async def get_all_links() -> list:
        """Get all links on the page with text and href."""
        await log_chat("Getting links")
        return await browser_command("get_links", {})

    @tool
    async def get_search_results() -> list:
        """Extract search results from Google, DuckDuckGo, or Brave. Returns [{title, url, snippet}]."""
        await log_chat("Extracting search results")
        return await browser_command("get_search_results", {})

    @tool
    async def get_all_headings() -> list:
        """Get all headings (H1-H6) on the page."""
        await log_chat("Getting headings")
        return await browser_command("get_headings", {})

    @tool
    async def get_ui_schema(mode: str = "visible") -> list:
        """Extract interactive elements from the page (buttons, inputs, links, etc.). Modes: visible (default), full."""
        await log_chat("Getting UI schema")
        return await browser_command("get_schema", {"mode": mode})

    @tool
    async def get_page_content() -> str:
        """Get the raw HTML of the page body."""
        return await browser_command("get_page_content", {})

    # ── USER INPUT ───────────────────────────────
    # These tools are interrupted by HumanInTheLoopMiddleware via interrupt_on.
    # The tool functions are essentially stubs — the user's respond value
    # replaces the tool output entirely (the tool never executes).

    @tool
    async def get_user_confirmation(query: str) -> str:
        """Ask user a yes/no question. Returns the user's response."""
        return query

    @tool
    async def get_user_input_from_options(options: List[str]) -> str:
        """Ask the user to choose one option from the given list. Pass options as a JSON list of strings, e.g. ["Red", "Green", "Blue"]. Returns the user's chosen option text."""
        if isinstance(options, str):
            try:
                options = json.loads(options)
            except json.JSONDecodeError:
                return "ERROR: options must be a valid JSON list of strings."
        if not isinstance(options, list) or not all(isinstance(o, str) for o in options):
            return "ERROR: options must be a valid JSON list of strings."
        return json.dumps(options)

    # ── FILE OPERATIONS ──────────────────────────

    @tool
    async def write_file(content: str, filename: str, append: bool = False) -> str:
        """Write content to a file. Set append=true to add to end (text files only).

        Supported formats (chosen by filename extension):
        - .txt/.md/.csv/.json and other text: written as-is; append supported.
        - .docx: markdown headings (#, ##, ###) become document headings and markdown tables (| a | b |) become tables.
        - .xlsx: each line becomes a row; cells separated by tabs (or commas).
        - .pptx: provide the slide deck as HTML+CSS. Each element with class="slide" (or a <section>) becomes one slide; the first <h1> is the slide title, the rest is the body. Inline style font-size/color are applied. Plain markdown (slides split by ---) is also accepted. For exact control over each slide, use the create_pptx tool instead (it accepts per-slide PresentationML XML).
        - .pdf: provide the document as a full HTML+CSS page (e.g. starting with <!DOCTYPE html> or <html>); it is rendered with a headless browser so all CSS applies. Plain markdown is also accepted as a fallback.
        """
        await log_chat(f"Writing {filename}")
        try:
            full_path = resolve_user_path(filename)
            os.makedirs(os.path.dirname(full_path), exist_ok=True)
            return await write_text_to_file(full_path, content, append)
        except Exception as e:
            return str(e)

    @tool
    async def create_pptx(filename: str, slides: str) -> str:
        """Create an editable PowerPoint (.pptx) where every slide is authored as raw PresentationML XML.

        `slides` is a JSON array — one object per slide:
        [{"title": "Q3 Results", "xml": "<p:txBody>...</p:txBody>"}, ...]

        Per-slide fields:
        - title: plain-text slide title (optional).
        - body: plain text (multi-line) or list of lines for the body placeholder; lines starting with "- " become bullets (optional).
        - xml: raw PresentationML fragment for full control (optional; takes priority over title/body).
        - mode: "txbody" (default) | "spTree" | "append".
        - placeholder: target placeholder for txbody mode — "body" (default), "title", any placeholder type, or "idx:<number>" to match by index.

        XML modes:
        - txbody: xml is a <p:txBody>...</p:txBody> fragment replacing the targeted placeholder's text.
        - spTree: xml is one or more <p:sp>...</p:sp> shapes replacing every shape on the slide (placeholders are removed).
        - append: xml is one or more <p:sp>...</p:sp> shapes added after the existing placeholders.

        XML rules: use only the standard prefixes p:, a:, r: and never add xmlns declarations; sizes are in hundredths of a point (1400 = 14pt, 3200 = 32pt); bold with b="1".
        """
        await log_chat(f"Creating presentation {filename}")
        try:
            specs = json.loads(slides)
        except Exception as e:
            return f"ERROR: slides is not valid JSON: {e}"
        if not isinstance(specs, list) or not specs:
            return "ERROR: slides must be a non-empty JSON array"
        try:
            from pptx_builder import PptxBuilder, title_txbody, body_txbody
            full_path = resolve_user_path(filename)
            if not full_path.lower().endswith(".pptx"):
                full_path += ".pptx"
            builder = PptxBuilder(num_slides=len(specs))
            try:
                for i, spec in enumerate(specs, start=1):
                    if not isinstance(spec, dict):
                        return f"ERROR: slide {i} must be a JSON object"
                    mode = str(spec.get("mode", "txbody")).lower()
                    xml = spec.get("xml")
                    placeholder_type, placeholder_idx = None, None
                    ph = spec.get("placeholder")
                    if ph is not None:
                        ph = str(ph)
                        if ph.startswith("idx:"):
                            placeholder_idx = ph[4:]
                        else:
                            placeholder_type = ph
                    if mode == "sptree":
                        if not xml:
                            return f"ERROR: slide {i} mode=spTree requires 'xml'"
                        builder.set_slide_shapes_xml(i, xml)
                    elif mode == "txbody":
                        if spec.get("title"):
                            builder.set_placeholder_xml(i, title_txbody(spec["title"]), placeholder_type="title")
                        if xml:
                            builder.set_placeholder_xml(i, xml, placeholder_type=placeholder_type, placeholder_idx=placeholder_idx)
                        elif spec.get("body"):
                            body = spec["body"] if isinstance(spec["body"], (list, tuple)) else str(spec["body"])
                            builder.set_placeholder_xml(i, body_txbody(body), placeholder_type=placeholder_type, placeholder_idx=placeholder_idx)
                    elif mode == "append":
                        if not xml:
                            return f"ERROR: slide {i} mode=append requires 'xml'"
                        if spec.get("title"):
                            builder.set_placeholder_xml(i, title_txbody(spec["title"]), placeholder_type="title")
                        if spec.get("body"):
                            body = spec["body"] if isinstance(spec["body"], (list, tuple)) else str(spec["body"])
                            builder.set_placeholder_xml(i, body_txbody(body), placeholder_type=placeholder_type, placeholder_idx=placeholder_idx)
                        builder.append_shapes_xml(i, xml)
                    else:
                        return f"ERROR: slide {i} unknown mode {mode!r}"
                builder.save(full_path)
            finally:
                builder.cleanup()
            return f"OK: {os.path.basename(full_path)}"
        except Exception as e:
            return f"ERROR: {e}"

    @tool
    async def read_file(filepath: str) -> str:
        """Read file contents. Supports PDF, DOCX, XLSX, PPTX, and plain text files."""
        await log_chat(f"Reading {filepath}")
        try:
            full_path = resolve_user_path(filepath)
            if not os.path.exists(full_path):
                return f"Not found: {filepath}"
            return extract_text_from_file(full_path)
        except Exception as e:
            return str(e)

    @tool
    async def delete_entry(filepath: str) -> str:
        """Delete a file or directory."""
        import shutil
        await log_chat(f"Deleting {filepath}")
        try:
            full_path = resolve_user_path(filepath)
            if not os.path.exists(full_path):
                return f"Not found: {filepath}"
            if os.path.isdir(full_path):
                shutil.rmtree(full_path)
            else:
                os.remove(full_path)
            return f"OK: {filepath}"
        except Exception as e:
            return str(e)

    # ── UTILITY ──────────────────────────────────

    @tool
    async def get_current_date_time() -> str:
        """Get current date and time."""
        return datetime.datetime.now().strftime("%A, %d %B %Y %H:%M:%S")

    @tool
    async def get_total_tokens(pattern: str) -> str:
        """Get token count for files matching a regex pattern in the project directory. Use before reading any file to check size. Estimates ~4 chars per token."""
        import re as _re
        total_chars = 0
        file_count = 0
        matched_files = []
        try:
            files_dir = get_user_files_dir()
            if not os.path.isdir(files_dir):
                return "No files directory found."
            regex = _re.compile(pattern, _re.IGNORECASE)
            for root, dirs, files in os.walk(files_dir):
                for fname in files:
                    if fname.startswith("."):
                        continue
                    if not regex.search(fname):
                        continue
                    fpath = os.path.join(root, fname)
                    try:
                        size = os.path.getsize(fpath)
                        total_chars += size
                        file_count += 1
                        rel = os.path.relpath(fpath, files_dir)
                        matched_files.append(f"{rel} ({size:,} chars)")
                    except Exception:
                        continue
            total_tokens = total_chars // 4
            if file_count == 0:
                return f"No files matching pattern '{pattern}'"
            listing = "\n".join(matched_files[:20])
            more = f"\n... and {file_count - 20} more" if file_count > 20 else ""
            return f"Matched: {file_count} files\n{listing}{more}\n\nEstimated tokens: ~{total_tokens:,}"
        except Exception as e:
            return f"Error: {e}"

    @tool
    async def action_logger(action: str) -> str:
        """Log an action to the user."""
        try:
            await log_chat(action)
            return "OK"
        except:
            return "Error"

    # ── TOOL LIST ────────────────────────────────

    return [
        open_url, get_url, get_title, go_back, go_forward,
        click, type_text, scroll, submit_form, press_key, page_down, page_up,
        select_option, get_dropdown_options,
        get_page_text, get_all_links, get_search_results, get_all_headings, get_ui_schema, get_page_content,
        get_user_confirmation, get_user_input_from_options,
        write_file, create_pptx, read_file, delete_entry,
        get_current_date_time, get_total_tokens, action_logger,
    ]
